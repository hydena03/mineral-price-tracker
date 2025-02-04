from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

def apply_slide_design(slide):
    """슬라이드에 디자인을 적용하는 함수"""
    # 배경 도형 추가
    left = 0
    top = 0
    width = Inches(10)
    height = Inches(7.5)
    
    # 상단 장식 도형 추가
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(1))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 141)  # 진한 파란색
    shape.shadow.inherit = False
    
    # 하단 장식 도형 추가
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(6.5), width, Inches(1))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 141)  # 진한 파란색
    shape.shadow.inherit = False

def apply_theme_colors(slide):
    """슬라이드에 테마 색상을 적용하는 함수"""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if shape.name == "Title":
                        run.font.color.rgb = RGBColor(255, 255, 255)  # 흰색
                        run.font.bold = True
                        run.font.size = Pt(32)
                    else:
                        run.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
                        run.font.size = Pt(18)

def create_title_slide(prs, title, subtitle):
    """타이틀 슬라이드를 생성하는 함수"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경 디자인 적용
    left = 0
    top = 0
    width = Inches(10)
    height = Inches(7.5)
    
    # 전체 배경 도형
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 141)  # 진한 파란색
    
    # 제목 추가
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    
    # 부제목 추가
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    
    return slide

def create_content_slide(prs, title, content, level=1):
    """내용 슬라이드를 생성하는 함수"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경 디자인 적용
    apply_slide_design(slide)
    
    # 제목 추가
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # 내용 추가
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # 여러 줄의 내용을 처리
    first = True
    for line in content.split('\n'):
        if first:
            p = content_frame.paragraphs[0]
            first = False
        else:
            p = content_frame.add_paragraph()
        
        p.text = line
        p.level = level if line.strip().startswith('•') else 0
        p.font.size = Pt(18)
        if line.strip().startswith('•'):
            p.font.color.rgb = RGBColor(0, 51, 141)  # 진한 파란색
        else:
            p.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
    
    return slide

def add_footer(slide, text="AMZ Bank © 2024"):
    """슬라이드에 푸터를 추가하는 함수"""
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(2)
    height = Inches(0.3)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색

def create_amz_bank_presentation():
    try:
        prs = Presentation()
        
        # 기본 슬라이드 크기 설정 (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 타이틀 슬라이드 생성
        title_slide = create_title_slide(
            prs,
            "AMZ Bank\n인터넷전문은행 인가 신청 프로세스",
            "디지털 뱅크의 새로운 시작"
        )
        
        # 진행 일정 슬라이드
        timeline_content = """주요 일정:

• 2024년 12월 12일~17일: 인가신청 희망사업자 의견수렴 기간
• 2025년 3월 25일~26일: 예비인가 신청서 접수 기간
• 2025년 5월 (잠정): 예비인가 심사결과 발표
• 2025년 하반기: 본인가 진행 (예비인가 취득 시)"""
        create_content_slide(prs, "진행 일정", timeline_content)

        # 필수 제출 서류 슬라이드
        docs_content = """필수 제출 서류:

• 예비인가 신청서 (3부)
• 사업계획 관련 서류
• 자본금 조달 및 주주구성 계획서
• IT 인프라 및 보안 계획서
• 리스크 관리체계 문서
• 소비자보호 방안
• 지배구조 관련 서류"""
        create_content_slide(prs, "필수 제출 서류", docs_content)

        # 진행 절차 슬라이드
        process_content = """진행 절차:

• 서류 준비 (2024년 12월 - 2025년 3월)
• 신청서 제출 (2025년 3월 25일-26일)
• 외부평가위원회 심사
• 금융감독원 심사
• 금융위원회 심의 및 의결"""
        create_content_slide(prs, "진행 절차", process_content)

        # 인가 획득을 위한 중점 고려사항 슬라이드들
        considerations = {
            "자금조달의 안정성 (150점)": """• 충분한 초기 자본금
• 안정적인 자금조달 방안
• 명확한 추가 자본조달 계획""",
            
            "사업계획의 혁신성 (350점)": """• 차별화된 금융서비스 제공
• 혁신적 신용평가모형 구축
• 금융과 ICT의 융합
• 시장 경쟁력 강화 가능성""",
            
            "포용성 (200점)": """• 금융소외계층 지원방안
• 중금리 대출 프로그램
• 지역금융 기여도
• 중소기업 지원 프로그램""",
            
            "실현가능성 및 리스크 관리 (300점)": """• 사업계획의 실현가능성
• 리스크 관리 체계
• IT 보안 체계
• 소비자보호 체계"""
        }

        for title, content in considerations.items():
            create_content_slide(prs, title, content)

        # 성공을 위한 핵심 요소 슬라이드들
        success_factors = {
            "기술 융합": """• 첨단 디지털 플랫폼 구축
• 혁신적 금융서비스 개발
• 안전한 IT 인프라 확보""",
            
            "시장 차별화": """• 독특한 가치 제안
• 명확한 목표 고객층 설정
• 경쟁력 있는 서비스 제공""",
            
            "규제 준수": """• 금융위원회 가이드라인 준수
• 리스크 관리 체계 구축
• 소비자보호 방안 마련"""
        }

        for title, content in success_factors.items():
            create_content_slide(prs, f"성공 요소: {title}", content)

        # 향후 준비 계획 슬라이드들
        preparation_plans = {
            "즉시 실행 과제 (2024년 12월 - 2025년 1월)": """• 인가 신청 전담팀 구성
• 서류 준비 착수
• 상세 사업계획 수립""",
            
            "신청 전 준비 단계 (2025년 2월 - 3월)": """• 내부 검토 및 보완
• 법률 준수 여부 확인
• 최종 서류 완성""",
            
            "신청 기간 (2025년 3월 25일-26일)": """• 완성된 신청 서류 제출
• 추가 질의 대비
• 보완 자료 준비"""
        }

        for title, content in preparation_plans.items():
            create_content_slide(prs, f"준비 계획: {title}", content)
        
        # 각 슬라이드에 푸터 추가
        for slide in prs.slides:
            add_footer(slide)
        
        prs.save('AMZ_Bank_인가신청_프로세스.pptx')
        print("프레젠테이션이 성공적으로 생성되었습니다.")
        
    except Exception as e:
        print(f"프레젠테이션 생성 중 오류가 발생했습니다: {str(e)}")

if __name__ == "__main__":
    create_amz_bank_presentation() 